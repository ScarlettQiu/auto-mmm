"""Generate markdown + PowerPoint analysis report from MMM results.

Usage:
    python report_builder.py                  # from latest results
    python report_builder.py --round 3
    python report_builder.py --summary "Claude's interpretation text"
"""
from __future__ import annotations

import argparse
import io
import json
from collections import Counter
from pathlib import Path

from compare import (
    contribution_comparison,
    disagreements,
    model_fit_summary,
    roi_comparison,
    top_channels,
)

# ── Model metadata ────────────────────────────────────────────────────────────

MODEL_META = {
    "ridge": {
        "full_name": "Ridge Regression MMM",
        "colour_hex": (0x1F, 0x77, 0xB4),
        "tagline": "Fast, transparent baseline — always runs",
        "how_it_works": (
            "Fits a regularised linear regression on adstock-transformed spend. "
            "Bootstraps 200 samples to estimate coefficient uncertainty. "
            "Alpha (regularisation strength) selected via 5-fold cross-validation."
        ),
        "pros": [
            "Runs in seconds — no heavy dependencies",
            "Coefficients are fully interpretable",
            "Bootstrap gives confidence intervals per channel",
            "Good baseline to sanity-check other models against",
        ],
        "cons": [
            "Assumes linear spend-response (no saturation curve)",
            "Prone to overfitting on small samples (R²=1.0 is a warning sign)",
            "Shrinks small coefficients to zero — may under-attribute niche channels",
            "No uncertainty propagation — point estimates only",
        ],
        "best_for": [
            "Quick first-pass analysis when data is limited",
            "Sanity-checking Bayesian results",
            "Teams without Python/stats expertise (simple to explain)",
            "Datasets with 50+ weekly observations and low collinearity",
        ],
        "avoid_when": [
            "Sample size is very small (<30 periods) — overfitting risk",
            "Channels have strong diminishing returns (Hill curve needed)",
            "You need full uncertainty quantification for budget decisions",
        ],
    },
    "pymc": {
        "full_name": "Bayesian MMM (PyMC-Marketing)",
        "colour_hex": (0xFF, 0x7F, 0x0E),
        "tagline": "Full posterior uncertainty — the gold standard",
        "how_it_works": (
            "Uses DelayedSaturatedMMM from PyMC-Marketing. Fits a full Bayesian model "
            "with prior distributions on adstock decay, Hill saturation, and channel "
            "coefficients. MCMC sampling (NUTS) produces a posterior distribution over "
            "all parameters — every ROI estimate comes with a credible interval."
        ),
        "pros": [
            "Full uncertainty quantification — know how confident you are",
            "Encodes domain knowledge via priors (e.g. TV decays slower than Digital)",
            "Handles small samples better by regularising through priors",
            "Posterior predictive checks catch model misspecification",
            "Industry standard for rigorous MMM (used by Meta, Google)",
        ],
        "cons": [
            "Slow — 10–60 minutes depending on data size and chain count",
            "Requires pymc-marketing installation and familiarity with Bayesian stats",
            "Prior choices matter — wrong priors can bias results",
            "Results harder to explain to non-technical stakeholders",
        ],
        "best_for": [
            "Production MMM where budget decisions involve significant spend",
            "Datasets with 80+ weekly observations",
            "When you need credible intervals on ROI for board-level decisions",
            "Teams with data science capability to interpret posteriors",
        ],
        "avoid_when": [
            "You need results in minutes (use Ridge first)",
            "Fewer than 30 observations — priors will dominate data",
            "No one on the team can interpret MCMC diagnostics",
        ],
    },
    "lightweight_mmm": {
        "full_name": "LightweightMMM / NNLS",
        "colour_hex": (0x2C, 0xA0, 0x2C),
        "tagline": "Google's JAX-based MMM with NNLS fallback",
        "how_it_works": (
            "Google's LightweightMMM uses JAX for fast GPU/CPU Bayesian inference with "
            "built-in Hill adstock. When JAX is unavailable, falls back to Non-Negative "
            "Least Squares (NNLS via scipy) — a constrained regression that enforces "
            "positive channel coefficients, preventing the negative-ROI problem common in Ridge."
        ),
        "pros": [
            "NNLS fallback enforces positive ROI — no negative coefficients",
            "LightweightMMM is significantly faster than PyMC on GPU",
            "Built-in Hill + adstock transforms — no manual feature engineering",
            "Open source from Google — well tested at scale",
            "Good middle ground between Ridge speed and PyMC rigour",
        ],
        "cons": [
            "NNLS fallback loses uncertainty quantification",
            "JAX installation can be complex (especially on Mac M-series)",
            "LightweightMMM is less actively maintained than PyMC-Marketing",
            "NNLS can over-attribute to channels that happen to correlate with KPI",
        ],
        "best_for": [
            "Teams who want faster Bayesian inference than PyMC",
            "Google Cloud / GCP environments where JAX runs natively",
            "As a cross-check against PyMC results",
            "When you need positive-constrained estimates without full Bayes",
        ],
        "avoid_when": [
            "Full uncertainty quantification is required (use PyMC instead)",
            "You cannot install JAX (NNLS fallback loses most benefits)",
            "Data has high collinearity — NNLS will over-attribute to one channel",
        ],
    },
}

MODEL_COMPARISON_TABLE = [
    # (Dimension, ridge, pymc, lightweight_mmm)
    ("Speed",               "⚡ Seconds",        "🐢 10–60 min",     "🔄 Minutes (JAX)"),
    ("Uncertainty",         "Bootstrap CIs",     "Full posterior",    "NNLS: none / JAX: posterior"),
    ("Saturation curve",    "❌ Linear only",    "✅ Hill (built-in)", "✅ Hill (built-in)"),
    ("Adstock",             "Manual (prepare.py)","✅ Built-in",      "✅ Built-in"),
    ("Negative ROI risk",   "⚠️ Yes (shrinkage)","Low (priors)",     "✅ No (NNLS positive)"),
    ("Min data needed",     "50+ weekly obs",    "80+ weekly obs",    "50+ weekly obs"),
    ("Install complexity",  "✅ pip sklearn",    "⚠️ pymc-marketing", "⚠️ JAX required"),
    ("Explainability",      "✅ High",           "Medium",            "Medium"),
    ("Industry adoption",   "Baseline standard", "⭐ Gold standard",  "Google internal"),
]

WHEN_TO_USE = [
    {
        "scenario": "Quick first-pass analysis",
        "data": "Any size",
        "recommended": "Ridge",
        "reason": "Results in seconds, easy to interpret, good sanity check",
    },
    {
        "scenario": "Small dataset (<30 periods)",
        "data": "Monthly, <30 obs",
        "recommended": "PyMC",
        "reason": "Priors regularise better than Ridge shrinkage on thin data",
    },
    {
        "scenario": "Production budget decisions",
        "data": "Weekly, 80+ obs",
        "recommended": "PyMC",
        "reason": "Full credible intervals needed for significant spend decisions",
    },
    {
        "scenario": "Fast iteration / experiment loop",
        "data": "Weekly, 50+ obs",
        "recommended": "Ridge → LightweightMMM",
        "reason": "Ridge for speed, LightweightMMM for positive-constrained check",
    },
    {
        "scenario": "GCP / GPU environment",
        "data": "Any size",
        "recommended": "LightweightMMM",
        "reason": "JAX runs natively on GCP — fastest Bayesian option",
    },
    {
        "scenario": "Board-level reporting",
        "data": "Weekly, 100+ obs",
        "recommended": "PyMC + Ridge as sanity check",
        "reason": "Posterior credible intervals defensible; Ridge validates direction",
    },
    {
        "scenario": "High collinearity between channels",
        "data": "Any",
        "recommended": "PyMC",
        "reason": "Priors on each channel help disentangle correlated spend patterns",
    },
]


# ── Exploration report loader ─────────────────────────────────────────────────

def load_exploration_report() -> str:
    """Return the content of rounds/R01_data_exploration.md, or '' if missing."""
    path = Path("rounds/R01_data_exploration.md")
    if path.exists():
        return path.read_text()
    return ""


def _parse_exploration(text: str) -> dict:
    """Extract readiness score and key section snippets from exploration report."""
    out = {"score": "", "sections": {}, "raw": text}
    for line in text.splitlines():
        if line.startswith("**Readiness Score:**") or line.startswith("**Score:"):
            out["score"] = line.strip()
        if line.startswith("## "):
            current = line[3:].strip()
            out["sections"][current] = []
        elif out["sections"]:
            last = list(out["sections"])[-1]
            out["sections"][last].append(line)
    # Trim trailing blank lines per section
    for k in out["sections"]:
        while out["sections"][k] and not out["sections"][k][-1].strip():
            out["sections"][k].pop()
    return out


# ── Markdown report ──────────────────────────────────────────────────────────

def build_markdown(results: dict, summary_text: str = "") -> str:
    roi_df     = roi_comparison(results)
    contrib_df = contribution_comparison(results)
    fit_df     = model_fit_summary(results)
    top        = top_channels(results)
    disc       = disagreements(roi_df)

    run_at = results.get("run_at", "")[:10]
    lines = [
        "# MMM Analysis Report",
        f"*Generated {run_at} · {results.get('train_periods', '?')} training periods · "
        f"{results.get('test_periods', '?')} holdout periods*",
        "",
    ]

    if summary_text:
        lines += ["## Executive Summary", "", summary_text, ""]

    # Data exploration findings
    exploration_text = load_exploration_report()
    if exploration_text:
        lines += ["## Data Exploration Findings", "", exploration_text, ""]

    # Model introductions
    lines += ["## The Three Models", ""]
    for key, m in MODEL_META.items():
        lines += [
            f"### {m['full_name']}",
            f"*{m['tagline']}*",
            "",
            m["how_it_works"],
            "",
            "**Pros:** " + " · ".join(m["pros"]),
            "",
            "**Cons:** " + " · ".join(m["cons"]),
            "",
            "**Best for:** " + " · ".join(m["best_for"]),
            "",
        ]

    lines += [
        "## Model Comparison",
        "",
        "| Dimension | Ridge | PyMC | LightweightMMM |",
        "|---|---|---|---|",
    ]
    for row in MODEL_COMPARISON_TABLE:
        lines.append(f"| {row[0]} | {row[1]} | {row[2]} | {row[3]} |")
    lines.append("")

    lines += ["## When to Use Which Model", ""]
    lines += ["| Scenario | Data | Recommended | Reason |", "|---|---|---|---|"]
    for w in WHEN_TO_USE:
        lines.append(f"| {w['scenario']} | {w['data']} | **{w['recommended']}** | {w['reason']} |")
    lines.append("")

    lines += [
        "## Model Fit (This Run)",
        "",
        fit_df.to_markdown() if hasattr(fit_df, "to_markdown") else fit_df.to_string(),
        "",
        "## ROI by Channel",
        "",
        "*Higher = more revenue per unit of spend. Agreement = cross-model consistency.*",
        "",
        roi_df.to_markdown() if hasattr(roi_df, "to_markdown") else roi_df.to_string(),
        "",
        "## Channel Contribution (%)",
        "",
        contrib_df.to_markdown() if hasattr(contrib_df, "to_markdown") else contrib_df.to_string(),
        "",
        "## Top Channels by Model",
        "",
    ]
    for model, chans in top.items():
        lines.append(f"- **{model}**: {', '.join(chans)}")
    lines.append("")

    if not disc.empty:
        lines += [
            "## High Disagreement Channels",
            "",
            "*Interpret with caution — models give significantly different answers.*",
            "",
            disc.to_markdown() if hasattr(disc, "to_markdown") else disc.to_string(),
            "",
        ]

    lines += [
        "## Data Caveat",
        "",
        f"This analysis uses {results.get('data_periods', '?')} monthly periods. "
        "MMM typically benefits from 2+ years of weekly data. "
        "Interpret uncertainty ranges with this in mind.",
        "",
    ]

    return "\n".join(lines)


# ── PowerPoint helpers ────────────────────────────────────────────────────────

def _make_pptx_helpers(prs, W, H, MARGIN, BLUE, DARK, GRAY, WHITE, MID, Inches, Pt):

    def blank():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def rect(slide, l, t, w, h, color):
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def text(slide, l, t, w, h, txt, size=12, bold=False, color=DARK, wrap=True, font=None):
        txb = slide.shapes.add_textbox(l, t, w, h)
        tf = txb.text_frame; tf.word_wrap = wrap
        p = tf.paragraphs[0]; p.text = txt
        run = p.runs[0]
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color
        if font: run.font.name = font
        return txb

    def textbox_multiline(slide, l, t, w, h, lines_data, default_size=10):
        """lines_data: list of (text, size, bold, color)"""
        from pptx.util import Pt
        txb = slide.shapes.add_textbox(l, t, w, h)
        tf = txb.text_frame; tf.word_wrap = True
        for i, (txt, size, bold, color) in enumerate(lines_data):
            p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
            p.text = txt
            if p.runs:
                run = p.runs[0]
                run.font.size = Pt(size); run.font.bold = bold
                run.font.color.rgb = color

    def header(slide, title, dark=False):
        rect(slide, 0, 0, W, H * 0.14, DARK if dark else BLUE)
        text(slide, MARGIN, H * 0.02, W - MARGIN, H * 0.10,
             title, size=22, bold=True, color=WHITE)

    def section_divider(slide, title, color):
        rect(slide, 0, 0, W, H, color)
        text(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
             title, size=36, bold=True, color=WHITE)

    return blank, rect, text, textbox_multiline, header, section_divider


# ── PowerPoint builder ────────────────────────────────────────────────────────

def build_pptx(results: dict, summary_text: str = "") -> io.BytesIO:
    try:
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches, Pt
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    W = Inches(13.33)
    H = Inches(7.5)
    MARGIN = Inches(0.6)
    BLUE   = RGBColor(0x1F, 0x77, 0xB4)
    ORANGE = RGBColor(0xFF, 0x7F, 0x0E)
    GREEN  = RGBColor(0x2C, 0xA0, 0x2C)
    DARK   = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY   = RGBColor(0xF0, 0xF4, 0xF8)
    LGRAY  = RGBColor(0xD8, 0xE4, 0xF0)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    MID    = RGBColor(0x8A, 0x9B, 0xB0)
    RED    = RGBColor(0xC0, 0x39, 0x2B)
    TEAL   = RGBColor(0x16, 0xA0, 0x85)

    MODEL_COLORS = {"ridge": BLUE, "pymc": ORANGE, "lightweight_mmm": GREEN}
    MODEL_DARK   = {
        "ridge":           RGBColor(0x10, 0x50, 0x80),
        "pymc":            RGBColor(0xC0, 0x55, 0x00),
        "lightweight_mmm": RGBColor(0x18, 0x60, 0x18),
    }

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank, rect, text, textbox_multiline, header, section_divider = \
        _make_pptx_helpers(prs, W, H, MARGIN, BLUE, DARK, GRAY, WHITE, MID, Inches, Pt)

    run_at   = results.get("run_at", "")[:10]
    n_models = sum(1 for r in results["models"].values() if not r.get("skipped"))

    # ═════════════════════════════════════════════════════════════════════════
    # SECTION 1 — COVER
    # ═════════════════════════════════════════════════════════════════════════

    # Slide 1 — Title
    s = blank()
    rect(s, 0, 0, Inches(0.4), H, BLUE)
    rect(s, 0, H * 0.72, W, H * 0.28, GRAY)
    text(s, Inches(0.85), Inches(1.6), Inches(11.5), Inches(1.6),
         "Marketing Mix Model", size=38, bold=True, color=DARK)
    text(s, Inches(0.85), Inches(3.1), Inches(11.5), Inches(0.8),
         "Analysis Report", size=38, bold=False, color=BLUE)
    text(s, Inches(0.85), Inches(4.3), Inches(11.5), Inches(0.5),
         f"{results.get('train_periods','?')} training periods  ·  "
         f"{n_models} models compared  ·  {run_at}",
         size=13, color=MID)
    text(s, Inches(0.85), Inches(5.6), Inches(11.5), Inches(0.35),
         "Models: Ridge Regression  ·  Bayesian (PyMC)  ·  LightweightMMM",
         size=11, color=MID)

    # Slide 2 — Agenda
    s = blank()
    header(s, "Agenda")
    agenda = [
        ("01", "Model Overview",        "What are Ridge, PyMC, and LightweightMMM?"),
        ("02", "Pros, Cons & Scenarios","When to use each model and why"),
        ("03", "Model Comparison",      "Head-to-head on key dimensions"),
        ("04", "This Run — Fit Metrics","How well did each model perform?"),
        ("05", "ROI by Channel",        "Which channels generated the most return?"),
        ("06", "Channel Contribution",  "How much of GMV did each channel drive?"),
        ("07", "Agreement Analysis",    "Where do models agree and disagree?"),
        ("08", "Findings & Actions",    "What to do next"),
    ]
    for i, (num, title, sub) in enumerate(agenda):
        col = i % 2
        row = i // 2
        l = MARGIN + col * Inches(6.5)
        t = Inches(1.3) + row * Inches(1.38)
        rect(s, l, t, Inches(0.45), Inches(0.8), BLUE)
        text(s, l, t + Inches(0.1), Inches(0.45), Inches(0.6),
             num, size=14, bold=True, color=WHITE)
        text(s, l + Inches(0.55), t + Inches(0.05), Inches(5.7), Inches(0.38),
             title, size=13, bold=True, color=DARK)
        text(s, l + Inches(0.55), t + Inches(0.42), Inches(5.7), Inches(0.35),
             sub, size=10, color=MID)

    # ═════════════════════════════════════════════════════════════════════════
    # SECTION 2 — MODEL OVERVIEWS (one slide per model)
    # ═════════════════════════════════════════════════════════════════════════

    # Section divider
    s = blank()
    section_divider(s, "Section 1\nModel Overview", DARK)
    text(s, Inches(1.5), Inches(4.3), Inches(10), Inches(0.5),
         "Ridge Regression  ·  Bayesian (PyMC)  ·  LightweightMMM",
         size=16, color=WHITE)

    for key in ["ridge", "pymc", "lightweight_mmm"]:
        m = MODEL_META[key]
        color = MODEL_COLORS[key]
        dark  = MODEL_DARK[key]

        s = blank()
        # Colour band top
        rect(s, 0, 0, W, Inches(1.6), color)
        rect(s, 0, Inches(1.6), W, Inches(0.06), dark)
        text(s, MARGIN, Inches(0.15), W - MARGIN * 2, Inches(0.65),
             m["full_name"], size=26, bold=True, color=WHITE)
        text(s, MARGIN, Inches(0.82), W - MARGIN * 2, Inches(0.55),
             m["tagline"], size=14, color=WHITE)

        # How it works
        text(s, MARGIN, Inches(1.8), Inches(8.5), Inches(0.35),
             "How it works", size=12, bold=True, color=dark)
        text(s, MARGIN, Inches(2.15), Inches(8.5), Inches(1.1),
             m["how_it_works"], size=10, color=DARK)

        # Pros column
        pros_left = MARGIN
        text(s, pros_left, Inches(3.45), Inches(4.0), Inches(0.35),
             "✅  Pros", size=12, bold=True, color=TEAL)
        for i, pro in enumerate(m["pros"]):
            text(s, pros_left + Inches(0.1), Inches(3.85) + i * Inches(0.52),
                 Inches(3.8), Inches(0.48), f"· {pro}", size=10, color=DARK)

        # Cons column
        cons_left = Inches(5.3)
        text(s, cons_left, Inches(3.45), Inches(4.0), Inches(0.35),
             "⚠️  Cons", size=12, bold=True, color=RED)
        for i, con in enumerate(m["cons"]):
            text(s, cons_left + Inches(0.1), Inches(3.85) + i * Inches(0.52),
                 Inches(3.8), Inches(0.48), f"· {con}", size=10, color=DARK)

        # Best for box
        rect(s, Inches(9.6), Inches(1.8), Inches(3.15), Inches(5.5), GRAY)
        text(s, Inches(9.75), Inches(1.9), Inches(2.9), Inches(0.38),
             "Best for", size=11, bold=True, color=dark)
        for i, bf in enumerate(m["best_for"]):
            text(s, Inches(9.75), Inches(2.35) + i * Inches(0.72),
                 Inches(2.9), Inches(0.65), f"· {bf}", size=9, color=DARK)

    # ═════════════════════════════════════════════════════════════════════════
    # SECTION 3 — COMPARISON
    # ═════════════════════════════════════════════════════════════════════════

    s = blank()
    section_divider(s, "Section 2\nModel Comparison", RGBColor(0x2C, 0x3E, 0x50))
    text(s, Inches(1.5), Inches(4.3), Inches(10), Inches(0.5),
         "Pros, cons, when to use each model", size=16, color=WHITE)

    # Slide — comparison table
    s = blank()
    header(s, "Model Comparison — Key Dimensions")

    n_rows = len(MODEL_COMPARISON_TABLE) + 1
    tbl = s.shapes.add_table(
        n_rows, 4, MARGIN, Inches(1.15),
        W - MARGIN * 2, Inches(6.1)
    ).table
    col_ws = [Inches(3.2), Inches(2.9), Inches(2.9), Inches(3.1)]
    for ci, cw in enumerate(col_ws):
        tbl.columns[ci].width = cw

    headers = ["Dimension", "Ridge", "PyMC (Bayesian)", "LightweightMMM"]
    hcolors = [DARK, BLUE, ORANGE, GREEN]
    for j, (h_txt, h_col) in enumerate(zip(headers, hcolors)):
        cell = tbl.cell(0, j)
        cell.text = h_txt
        cell.fill.solid(); cell.fill.fore_color.rgb = h_col
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True; run.font.size = Pt(10)
                run.font.color.rgb = WHITE

    for i, row_data in enumerate(MODEL_COMPARISON_TABLE):
        bg = GRAY if i % 2 == 0 else WHITE
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)
                    run.font.color.rgb = BLUE if j == 0 else DARK
                    run.font.bold = (j == 0)

    # Slide — when to use
    s = blank()
    header(s, "When to Use Which Model")

    n_rows2 = len(WHEN_TO_USE) + 1
    tbl2 = s.shapes.add_table(
        n_rows2, 4, MARGIN, Inches(1.15),
        W - MARGIN * 2, Inches(6.1)
    ).table
    col_ws2 = [Inches(3.0), Inches(1.8), Inches(2.5), Inches(5.3)]
    for ci, cw in enumerate(col_ws2):
        tbl2.columns[ci].width = cw

    for j, (h_txt, h_col) in enumerate(zip(
        ["Scenario", "Data", "Recommended", "Why"],
        [DARK, DARK, DARK, DARK]
    )):
        cell = tbl2.cell(0, j)
        cell.text = h_txt
        cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True; run.font.size = Pt(10)
                run.font.color.rgb = WHITE

    rec_colors = {
        "Ridge": BLUE, "PyMC": ORANGE,
        "LightweightMMM": GREEN,
        "Ridge → LightweightMMM": TEAL,
        "PyMC + Ridge as sanity check": ORANGE,
    }
    for i, w in enumerate(WHEN_TO_USE):
        bg = GRAY if i % 2 == 0 else WHITE
        vals = [w["scenario"], w["data"], w["recommended"], w["reason"]]
        for j, val in enumerate(vals):
            cell = tbl2.cell(i + 1, j)
            cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)
                    run.font.color.rgb = rec_colors.get(val, DARK) if j == 2 else DARK
                    run.font.bold = (j == 2)

    # ═════════════════════════════════════════════════════════════════════════
    # SECTION 3b — DATA EXPLORATION (if available)
    # ═════════════════════════════════════════════════════════════════════════

    exploration_text = load_exploration_report()
    if exploration_text:
        expl = _parse_exploration(exploration_text)

        s = blank()
        section_divider(s, "Section 3\nData Exploration", RGBColor(0x14, 0x6C, 0x43))
        text(s, Inches(1.5), Inches(4.3), Inches(10), Inches(0.5),
             "Dataset quality, collinearity, anomalies, readiness score",
             size=16, color=WHITE)

        # Slide — exploration summary
        s = blank()
        header(s, "Data Exploration — Dataset Readiness")

        # Readiness score prominent display
        score_text = expl["score"] or "Readiness score not found"
        rect(s, MARGIN, Inches(1.25), Inches(4.2), Inches(1.2), TEAL)
        text(s, MARGIN + Inches(0.15), Inches(1.35), Inches(3.9), Inches(1.0),
             score_text, size=18, bold=True, color=WHITE)

        # Key sections on the right
        section_keys = ["Readiness Verdict", "7. Readiness Verdict",
                        "Anomalies", "5. Anomalies",
                        "Collinearity Check", "4. Collinearity Check"]
        y = Inches(1.25)
        for sk in section_keys:
            if sk in expl["sections"] and expl["sections"][sk]:
                snippet = "\n".join(expl["sections"][sk][:6]).strip()
                if snippet:
                    text(s, Inches(5.0), y, Inches(7.7), Inches(1.1),
                         f"{sk.split('. ')[-1]}: {snippet[:220]}", size=9, color=DARK)
                    y += Inches(1.15)
                    if y > Inches(6.5):
                        break

        # Full text in smaller font below if space allows
        flat_text = exploration_text[:900]
        text(s, MARGIN, Inches(2.6), W - MARGIN * 2, Inches(4.5),
             flat_text, size=8, color=DARK)

    # ═════════════════════════════════════════════════════════════════════════
    # SECTION 4 — THIS RUN'S RESULTS
    # ═════════════════════════════════════════════════════════════════════════

    s = blank()
    section_divider(s, "Section 4\nThis Run — Results", RGBColor(0x1A, 0x53, 0x76))
    text(s, Inches(1.5), Inches(4.3), Inches(10), Inches(0.5),
         "Model fit · ROI · Channel contribution · Agreement", size=16, color=WHITE)

    # Slide — model fit
    s = blank()
    header(s, "Model Fit — This Run")
    fit_df = model_fit_summary(results)

    card_w = Inches(3.8)
    card_h = Inches(2.5)
    x_positions = [MARGIN, Inches(4.85), Inches(9.1)]
    for i, (model_name, row) in enumerate(fit_df.iterrows()):
        if i >= 3: break
        color = MODEL_COLORS.get(model_name, BLUE)
        l = x_positions[i]
        rect(s, l, Inches(1.2), card_w, Inches(0.08), color)
        rect(s, l, Inches(1.28), card_w, card_h - Inches(0.08), GRAY)
        text(s, l + Inches(0.15), Inches(1.35), card_w - Inches(0.3), Inches(0.42),
             model_name.replace("_", " ").title(), size=13, bold=True, color=color)
        metrics = [
            ("Train R²",   str(row["train_r2"])),
            ("Train MAPE", str(row["train_mape"])),
            ("Test MAPE",  str(row["test_mape"])),
            ("Status",     str(row["status"])[:45]),
        ]
        for mi, (label, val) in enumerate(metrics):
            text(s, l + Inches(0.15), Inches(1.85) + mi * Inches(0.46),
                 Inches(1.3), Inches(0.38), label, size=9, color=MID)
            text(s, l + Inches(1.5), Inches(1.85) + mi * Inches(0.46),
                 card_w - Inches(1.65), Inches(0.38), val, size=9, bold=True, color=DARK)

    if summary_text:
        text(s, MARGIN, Inches(4.0), W - MARGIN * 2, Inches(3.2),
             summary_text[:500], size=10, color=DARK)

    # Slide — ROI chart
    roi_df = roi_comparison(results)
    if not roi_df.empty:
        s = blank()
        header(s, "ROI by Channel — All Models")
        model_cols = [c for c in roi_df.columns if c in MODEL_COLORS]
        channels = roi_df.index.tolist()

        chart_data = ChartData()
        chart_data.categories = channels
        for m in model_cols:
            chart_data.add_series(m.replace("_", " ").title(),
                                  [round(float(roi_df.loc[ch, m]), 4) for ch in channels])

        cs = s.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            MARGIN, Inches(1.15), W - MARGIN * 2, Inches(5.5), chart_data
        )
        chart = cs.chart
        chart.has_legend = True
        chart.legend.position = 2
        for i, m in enumerate(model_cols):
            chart.series[i].format.fill.solid()
            chart.series[i].format.fill.fore_color.rgb = MODEL_COLORS.get(m, BLUE)

        text(s, MARGIN, Inches(6.8), W - MARGIN * 2, Inches(0.4),
             "Higher = more revenue generated per unit of spend. "
             "Channels at zero may indicate insufficient data, not zero effect.",
             size=9, color=MID)

    # Slide — Contribution % chart
    contrib_df = contribution_comparison(results)
    if not contrib_df.empty:
        s = blank()
        header(s, "Channel Contribution % of GMV — All Models")
        model_cols = [c for c in contrib_df.columns if c in MODEL_COLORS]
        channels = contrib_df.index.tolist()

        chart_data = ChartData()
        chart_data.categories = channels
        for m in model_cols:
            chart_data.add_series(m.replace("_", " ").title(),
                                  [round(float(contrib_df.loc[ch, m]), 2) for ch in channels])

        cs = s.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            MARGIN, Inches(1.15), W - MARGIN * 2, Inches(5.5), chart_data
        )
        chart = cs.chart
        chart.has_legend = True
        chart.legend.position = 2
        for i, m in enumerate(model_cols):
            chart.series[i].format.fill.solid()
            chart.series[i].format.fill.fore_color.rgb = MODEL_COLORS.get(m, BLUE)

        text(s, MARGIN, Inches(6.8), W - MARGIN * 2, Inches(0.4),
             "Typical MMM: media explains 20–60% of KPI. "
             "Values near 0% suggest collinearity or insufficient data.",
             size=9, color=MID)

    # Slide — Agreement table
    s = blank()
    header(s, "Cross-Model Agreement by Channel")

    if not roi_df.empty:
        agree_df = roi_df[["mean_roi", "cv_pct", "agreement"]].reset_index()
        n = len(agree_df) + 1
        tbl3 = s.shapes.add_table(
            n, 4, MARGIN, Inches(1.2),
            W - MARGIN * 2, Inches(min(5.6, n * 0.6))
        ).table
        for ci, cw in enumerate([Inches(3.5), Inches(2.5), Inches(2.0), Inches(3.6)]):
            tbl3.columns[ci].width = cw

        for j, col in enumerate(["Channel", "Mean ROI", "CV %", "Agreement"]):
            cell = tbl3.cell(0, j)
            cell.text = col
            cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = True; run.font.size = Pt(10)
                    run.font.color.rgb = WHITE

        for i, (_, row) in enumerate(agree_df.iterrows()):
            bg = GRAY if i % 2 == 0 else WHITE
            cv = row["cv_pct"]
            agree_color = TEAL if cv < 20 else (ORANGE if cv < 50 else RED)
            vals = [row["channel"], f"{row['mean_roi']:.4f}", f"{cv:.1f}%", row["agreement"]]
            for j, val in enumerate(vals):
                cell = tbl3.cell(i + 1, j)
                cell.text = str(val)
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(9)
                        run.font.color.rgb = agree_color if j == 3 else DARK

        text(s, MARGIN, Inches(6.8), W - MARGIN * 2, Inches(0.45),
             "CV % = coefficient of variation across models. "
             "✅ <20% = high agreement  ⚠️ 20–50% = medium  ❌ >50% = low — interpret with caution",
             size=9, color=MID)

    # ═════════════════════════════════════════════════════════════════════════
    # SECTION 5 — FINDINGS & ACTIONS
    # ═════════════════════════════════════════════════════════════════════════

    s = blank()
    header(s, "Key Findings & Recommended Actions")

    top = top_channels(results)
    all_top = [ch for chans in top.values() for ch in chans]
    consensus = [ch for ch, cnt in Counter(all_top).most_common() if cnt >= 2]
    disc = disagreements(roi_df)

    bullet_y = Inches(1.25)
    text(s, MARGIN, bullet_y, W - MARGIN * 2, Inches(0.38),
         "Consensus top channels (appear in top 3 for 2+ models):",
         size=12, bold=True, color=DARK)
    bullet_y += Inches(0.45)
    text(s, MARGIN + Inches(0.2), bullet_y, W - MARGIN * 2, Inches(0.38),
         "  ·  ".join(consensus) if consensus else "No clear consensus — models disagree",
         size=12, color=BLUE if consensus else RED)
    bullet_y += Inches(0.65)

    if not disc.empty:
        text(s, MARGIN, bullet_y, W - MARGIN * 2, Inches(0.38),
             "⚠️  High disagreement channels (CV > 30%) — treat with caution:",
             size=12, bold=True, color=RED)
        bullet_y += Inches(0.45)
        text(s, MARGIN + Inches(0.2), bullet_y, W - MARGIN * 2, Inches(0.38),
             "  ·  ".join(disc.index.tolist()), size=12, color=RED)
        bullet_y += Inches(0.65)

    text(s, MARGIN, bullet_y, W - MARGIN * 2, Inches(0.38),
         "Recommended next steps:", size=12, bold=True, color=DARK)
    bullet_y += Inches(0.45)
    next_steps = [
        "Install PyMC-Marketing for full Bayesian uncertainty: pip install pymc-marketing",
        "Collect weekly data for 12+ months to reach minimum MMM threshold (52+ obs)",
        "Investigate any data anomalies before re-running (check for outlier periods)",
        "Run holdout test on top channel to validate model ROI estimate in practice",
    ]
    for step in next_steps:
        text(s, MARGIN + Inches(0.2), bullet_y, W - MARGIN * 2, Inches(0.42),
             f"→  {step}", size=10, color=DARK)
        bullet_y += Inches(0.5)

    rect(s, MARGIN, Inches(6.7), W - MARGIN * 2, Inches(0.5), LGRAY)
    text(s, MARGIN + Inches(0.15), Inches(6.75), W - MARGIN * 2 - Inches(0.3), Inches(0.4),
         f"Data caveat: {results.get('data_periods','?')} monthly periods used. "
         "MMM standard: 100+ weekly observations. Results are directional only.",
         size=9, color=MID)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── Entry point ───────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--round", type=int, default=None)
    parser.add_argument("--summary", default="")
    parser.add_argument("--config", default="config.json")
    parser.add_argument("--no-pptx", action="store_true")
    args = parser.parse_args()

    cfg = json.loads(Path(args.config).read_text())
    out_dir = Path(cfg.get("results_dir", "./results"))
    out_dir.mkdir(exist_ok=True)

    results_path = (
        Path("rounds") / f"R{args.round:02d}_results.json"
        if args.round else out_dir / "latest.json"
    )
    results = json.loads(results_path.read_text())

    md = build_markdown(results, args.summary)
    md_path = out_dir / "report.md"
    md_path.write_text(md)
    print(f"Markdown report → {md_path}")

    if not args.no_pptx:
        try:
            buf = build_pptx(results, args.summary)
            pptx_path = out_dir / "report.pptx"
            pptx_path.write_bytes(buf.read())
            print(f"PowerPoint deck → {pptx_path}")
        except ImportError as e:
            print(f"Skipping pptx: {e}")


if __name__ == "__main__":
    main()
